# -*- coding: utf-8 -*-
"""selftest.py — prove pptx_kit works without a real template.

Exercises the part that actually bit us: writing a speaker note to a notes slide whose master has
NO body placeholder (the case where python-pptx returns notes_text_frame is None). We reproduce that
by stripping the placeholder, then assert the note round-trips (write -> save -> reopen -> read).
    python selftest.py   # prints PASS / raises on failure
"""
import sys, io, os, tempfile
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
if hasattr(sys.stdout, "buffer"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx_kit import (new_deck, blank_slide_layout, speaker_note, fit_picture, overflows,
                      hang, text_units, wrapped_row_count, check_surface_leaks, save_and_check)

def _strip_notes_placeholders(slide):
    ns = slide.notes_slide                       # creates the notes slide
    for ph in list(ns.placeholders):
        ph._element.getparent().remove(ph._element)   # force notes_text_frame -> None

def test_text_and_leak_mechanics():
    # text_units: Hangul/Jamo count double, Latin counts single
    assert text_units("ab") == 2
    assert text_units("가나") == 4
    assert text_units("a가") == 3

    # wrapped_row_count: a line under wrap_units is 1 row; over it wraps to more
    assert wrapped_row_count(["short"], wrap_units=50) == 1
    assert wrapped_row_count(["가" * 30], wrap_units=50) == 2   # 30 hangul = 60 units -> wraps once
    assert wrapped_row_count(["a", "b"], wrap_units=50) == 2     # two separate 1-row lines

    # check_surface_leaks: no default term list (caller-supplied only), finds an exact hit
    prs = new_deck()
    s = prs.slides.add_slide(blank_slide_layout(prs))
    s.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "이전 버전과 비교하면"
    assert check_surface_leaks(prs, []) == [], "empty term list must report nothing"
    hits = check_surface_leaks(prs, ["이전 버전"])
    assert len(hits) == 1 and hits[0][1] == "이전 버전", f"expected one leak hit, got {hits}"

    # save_and_check: raises on a leak even though the file is on disk
    with tempfile.TemporaryDirectory() as d:
        out = os.path.join(d, "leak.pptx")
        try:
            save_and_check(prs, out, leak_terms=["이전 버전"])
            raise AssertionError("save_and_check should have raised on a known leak term")
        except SystemExit:
            pass
        assert os.path.exists(out), "save_and_check must still write the file before gating"

    print("PASS  text_units/wrapped_row_count estimate wrap correctly; "
          "check_surface_leaks/save_and_check gate on caller-supplied terms only.")

def main():
    prs = new_deck()                             # default template, blank deck
    assert len(prs.slides) == 0, "new_deck did not clear slides"
    s = prs.slides.add_slide(blank_slide_layout(prs))

    # force the hard path: no notes body placeholder, then write a multi-line note
    _strip_notes_placeholders(s)
    assert s.notes_slide.notes_text_frame is None, "expected a placeholder-less notes master"
    NOTE = "line one\nline two with an em-dash — and 한글"
    speaker_note(s, NOTE)

    # overflow helper sanity: a shape pushed off the slide is reported
    from pptx.util import Inches
    s.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(3), Inches(2))  # runs off 13.33x7.5
    assert overflows(s), "overflows() failed to flag an off-slide shape"

    # hang(): sets marL/indent on the paragraph's XML (python-pptx has no property for this)
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = "wraps under the bullet marker, not back to the margin"
    hang(p, 137160)  # ~0.15in in EMU
    assert p._pPr is not None and p._pPr.get("marL") == "137160" and p._pPr.get("indent") == "-137160", \
        "hang() did not set marL/indent as expected"

    with tempfile.TemporaryDirectory() as d:
        out = os.path.join(d, "selftest.pptx")
        prs.save(out)
        r = Presentation(out)
        got = r.slides[0].notes_slide.notes_text_frame.text
        assert got == NOTE, f"note round-trip mismatch:\n  wrote={NOTE!r}\n  read ={got!r}"

    print("PASS  speaker_note round-trips on a placeholder-less notes master; "
          "overflows() and hang() work.")

if __name__ == "__main__":
    main()
    test_text_and_leak_mechanics()
