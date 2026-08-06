# -*- coding: utf-8 -*-
"""pptx_kit — palette-agnostic PowerPoint MECHANICS, shared across the lab's deck builders.

Single home for the fiddly, template-dependent, easy-to-get-wrong parts of building a deck with
python-pptx. Project-specific STYLE (palette, fonts, card/section components) lives in each project's
own kit_common.py / ppt_common.py, which import the mechanics from here.

Why: speaker notes, template loading + slide clearing, blank-layout lookup, and overflow-safe image
placement are identical across every deck and each has a non-obvious failure mode. Duplicating them
per project is how they drift (e.g. one project "discovers" a template has no notes placeholder and
concludes notes are impossible — the fix is to inject the placeholder; see speaker_note).

Nothing here references a colour or font: callers pass those in. Keep it that way so it stays shared.
Import path: this file lives in the `pptx-editing` skill's scripts/; add that dir to sys.path.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def new_deck(template=None, width_in=13.333, height_in=7.5):
    """Load a template (inheriting its theme/master) and clear every slide, returning the deck.

    Removing the template's own slides while keeping its masters is the standard 'use the theme, build
    fresh' idiom, done by hand because python-pptx has no delete-slide API. slide_width/height are set;
    the caller tracks them if its styled components need SW/SH."""
    prs = Presentation(template) if template else Presentation()
    if not template:
        prs.slide_width = Inches(width_in); prs.slide_height = Inches(height_in)
    lst = prs.slides._sldIdLst
    for sid in list(lst):
        rId = (sid.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
               or sid.get("r:id"))
        if rId:
            try: prs.part.drop_rel(rId)
            except KeyError: pass
        lst.remove(sid)
    return prs


def blank_slide_layout(prs):
    """The blank layout. Templates differ: some expose a single placeholder-free 'DEFAULT' layout
    rather than the canonical index 6, so pick the first layout with no placeholders."""
    for lay in prs.slide_layouts:
        if len(lay.placeholders) == 0:
            return lay
    return prs.slide_layouts[0]


def rect(slide, l, t, w, h, color):
    """A filled, borderless rectangle (bars, panels, backgrounds). Coordinates in inches."""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background()
    return r


def slide_number(prs, slide, color, font, size=11, l=12.3, t=6.98, w=0.85, h=0.32):
    """Bottom-right slide number = current slide count. Colour and font are passed by the caller."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]; p.text = str(len(prs.slides))
    p.alignment = PP_ALIGN.RIGHT; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.name = font
    return box


def fit_picture(slide, path, left, top, max_w, max_h, center_x=True):
    """Add an image fitted inside a (max_w x max_h) box preserving aspect ratio, measured from the
    real pixel dimensions (PIL). Returns the Picture, or None if the file is missing (caller decides
    how to signal that). Overflow-safe: the image never exceeds the box. Coordinates in inches."""
    import os
    from PIL import Image
    if not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    x = left + (max_w - w) / 2 if center_x else left
    return slide.shapes.add_picture(path, Inches(x), Inches(top), Inches(w), Inches(h))


def overflows(slide, sw=13.333, sh=7.5, tol=0.02):
    """Return [(name, right_in, bottom_in), ...] for shapes whose right/bottom leaves the slide.
    Run before delivering: the list must be empty. python-pptx will not tell you otherwise."""
    out = []
    for sh_ in slide.shapes:
        try:
            r = (sh_.left + sh_.width) / 914400.0
            b = (sh_.top + sh_.height) / 914400.0
        except TypeError:
            continue  # some shapes (e.g. placeholders) may lack explicit geometry
        if r > sw + tol or b > sh + tol:
            out.append((sh_.name, round(r, 2), round(b, 2)))
    return out


def speaker_note(slide, text):
    """Attach a PowerPoint speaker note that SURVIVES a rebuild (it is written from the build script).

    Failure mode this guards against: some templates ship a notes master with no body placeholder, so
    python-pptx returns notes_text_frame is None and `slide.notes_slide.notes_text_frame.text = ...`
    raises. That is NOT 'this template can't hold notes' — PowerPoint shows the note fine; the
    placeholder is simply absent from the master. We inject a body placeholder into the notes slide's
    spTree directly, then the note renders normally. Newlines become separate paragraphs.

    A note written this way (from the script) is regenerated on every rebuild; a note typed into the
    .pptx by hand is wiped by the next rebuild. Put speaker notes in the script, not in PowerPoint."""
    from lxml import etree
    from xml.sax.saxutils import escape
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    if tf is not None:
        tf.text = text
        return
    paras = "".join("<a:p><a:r><a:t>%s</a:t></a:r></a:p>" % escape(ln) for ln in text.split("\n"))
    xml = ('<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
           '<p:nvSpPr><p:cNvPr id="10" name="Notes Placeholder"/>'
           '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
           '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
           '<p:spPr/><p:txBody><a:bodyPr/>' + paras + '</p:txBody></p:sp>')
    ns._element.spTree.append(etree.fromstring(xml))


# ---------------------------------------------------------------------------
# Native OOXML Math (mc:AlternateContent) — building NEW equations from scratch
# ---------------------------------------------------------------------------
# python-pptx has no math API at all, and (verified 2026-08-05 by introspecting the
# generated typelib module) common PowerPoint COM automation does NOT either: TextRange2
# exposes a `MathZones(Start, Length)` METHOD that returns an existing math zone as a
# TextRange2 — it reads, it does not create. There is no OMaths/MathZones.Add, no
# BuildUp, no LinearFormat setter anywhere in the typelib. Don't spend time on COM for
# this; SendKeys-driving the UI is the only COM-adjacent option and it's not worth it.
#
# What DOES work: place an empty textbox placeholder while building the deck, then AFTER
# prs.save() reopen the .pptx as a zip and string-replace that placeholder's <p:sp> with a
# real <mc:AlternateContent> block (Choice = native m:oMath, Fallback = plain text — the
# Fallback does NOT need to be a rendered image; PowerPoint only shows it to viewers that
# don't support the a14 math extension, so plain text is a perfectly valid fallback and
# skips the "render an image" problem entirely).
_MATH_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_MATH_RPR = ('<a:rPr><a:latin typeface="Cambria Math" panose="02040503050406030204" '
             'pitchFamily="18" charset="0"/></a:rPr>')
_MATH_CTRLPR = (f'<m:ctrlPr><a:rPr i="1"><a:latin typeface="Cambria Math" '
                f'panose="02040503050406030204" pitchFamily="18" charset="0"/></a:rPr></m:ctrlPr>')


def m_run(t):
    """A plain run inside a formula (one operator/variable/number)."""
    return f'<m:r>{_MATH_RPR}<m:t>{t}</m:t></m:r>'


def m_frac(num_xml, den_xml):
    """A fraction. num_xml/den_xml are fragments built from m_run/m_sub/m_sup/etc."""
    return f'<m:f><m:fPr>{_MATH_CTRLPR}</m:fPr><m:num>{num_xml}</m:num><m:den>{den_xml}</m:den></m:f>'


def m_sub(base_t, sub_t):
    """Subscript, e.g. n_rel."""
    return (f'<m:sSub><m:sSubPr>{_MATH_CTRLPR}</m:sSubPr>'
            f'<m:e>{m_run(base_t)}</m:e><m:sub>{m_run(sub_t)}</m:sub></m:sSub>')


def m_sup(base_xml, sup_t):
    """Superscript. Pass an m_sub(...) as base_xml to get a combined sub+superscript."""
    return (f'<m:sSup><m:sSupPr>{_MATH_CTRLPR}</m:sSupPr>'
            f'<m:e>{base_xml}</m:e><m:sup>{m_run(sup_t)}</m:sup></m:sSup>')


def m_nary(e_xml, chr_="∑"):
    """N-ary operator (Σ, Π, ∫, ...) with limits hidden (bare-symbol style)."""
    return (f'<m:nary><m:naryPr><m:chr m:val="{chr_}"/><m:limLoc m:val="undOvr"/>'
            f'<m:subHide m:val="on"/><m:supHide m:val="on"/>{_MATH_CTRLPR}</m:naryPr>'
            f'<m:sub/><m:sup/><m:e>{e_xml}</m:e></m:nary>')


def m_acc(base_t, chr_="̄"):
    """Accent over a single variable, e.g. x̄ (x-bar): m_acc('x'). Default chr_ is combining
    macron (U+0304, bar-above); pass another combining mark (e.g. '̂' for hat) for others."""
    return (f'<m:acc><m:accPr><m:chr m:val="{chr_}"/>{_MATH_CTRLPR}</m:accPr>'
            f'<m:e>{m_run(base_t)}</m:e></m:acc>')


def m_sqrt(e_xml):
    """Square root (degree hidden, i.e. a plain √, not an nth-root)."""
    return (f'<m:rad><m:radPr><m:degHide m:val="on"/>{_MATH_CTRLPR}</m:radPr>'
            f'<m:deg/><m:e>{e_xml}</m:e></m:rad>')


def equation_slot(slide, l, t, w, h, marker):
    """Empty placeholder textbox for a not-yet-real equation. promote_equations() finds it
    by this name (`EQN::marker`) after save and replaces it with the real oMath."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.name = f"EQN::{marker}"
    return box


def promote_equations(pptx_path, equations):
    """Call AFTER prs.save(pptx_path). equations = {marker: omath_run_xml, ...} where
    omath_run_xml is the m_run/m_frac/... fragments concatenated (no outer <m:oMath> tag —
    this function adds it). Rewrites each EQN::marker placeholder in the saved file into a
    real mc:AlternateContent equation, in place, preserving its position/size. Raises if a
    marker isn't found anywhere (fail loud, not silent).

    Gotcha this avoids: a naive extraction regex `<m:oMath[^>]*>` also matches the WRAPPER
    tag `<m:oMathPara ...>` (since `[^>]*` doesn't respect tag-name boundaries), silently
    nesting an extra oMathPara/oMath pair inside your formula. If you're lifting real m:oMath
    XML out of an existing hand-edited deck to reuse verbatim (rather than building fresh with
    m_run/m_frac/...), anchor the opening tag with a lookahead so it can't swallow "Para":
    `re.search(r'<m:oMath(?=[\\s>])[^>]*>(.*)</m:oMath>', xml, re.DOTALL)`.
    """
    import zipfile
    import re
    with zipfile.ZipFile(pptx_path, "r") as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    slide_names = [n for n in names if re.match(r"ppt/slides/slide\d+\.xml$", n)]
    remaining = dict(equations)
    for slide_name in slide_names:
        xml = data[slide_name].decode("utf-8")
        changed = False
        for marker in list(remaining):
            pat = re.compile(
                r'<p:sp><p:nvSpPr><p:cNvPr id="(\d+)" name="EQN::' + re.escape(marker) +
                r'"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>.*?</p:sp>', re.DOTALL)
            m = pat.search(xml)
            if not m:
                continue
            shape_id = m.group(1)
            xfrm_m = re.search(r'<a:off x="(-?\d+)" y="(-?\d+)"/><a:ext cx="(\d+)" cy="(\d+)"/>', m.group(0))
            off_x, off_y, ext_cx, ext_cy = xfrm_m.groups()
            omath = f'<m:oMath xmlns:m="{_MATH_NS}">{remaining[marker]}</m:oMath>'
            choice_sp = (
                f'<p:sp><p:nvSpPr><p:cNvPr id="{shape_id}" name="EQN::{marker}"/>'
                '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>'
                '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
                '<p:txBody><a:bodyPr wrap="square"><a:spAutoFit/></a:bodyPr><a:lstStyle/>'
                f'<a:p><a:pPr algn="l"/><a14:m><m:oMathPara xmlns:m="{_MATH_NS}">'
                f'<m:oMathParaPr><m:jc m:val="left"/></m:oMathParaPr>{omath}</m:oMathPara></a14:m>'
                '<a:endParaRPr/></a:p></p:txBody></p:sp>'
            )
            fallback_sp = (
                f'<p:sp><p:nvSpPr><p:cNvPr id="{shape_id}" name="EQN::{marker}"/>'
                '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
                f'<p:spPr><a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>'
                '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/></p:spPr>'
                '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:r><a:rPr lang="ko-KR" altLang="en-US" sz="1800"/>'
                f'<a:t>[{marker}]</a:t></a:r></a:p></p:txBody></p:sp>'
            )
            new_shape = (
                '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
                'xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main">'
                f'<mc:Choice Requires="a14">{choice_sp}</mc:Choice>'
                f'<mc:Fallback xmlns="">{fallback_sp}</mc:Fallback></mc:AlternateContent>'
            )
            xml = xml[:m.start()] + new_shape + xml[m.end():]
            changed = True
            del remaining[marker]
        if changed:
            data[slide_name] = xml.encode("utf-8")

    if remaining:
        raise SystemExit(f"promote_equations: markers not found: {list(remaining)} — "
                          f"check equation_slot() names/slides")

    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])


def hang(paragraph, width):
    """Hanging indent of `width` (an EMU-ish int, e.g. Pt(sz)*0.95) on one paragraph.

    python-pptx exposes no API for this. Without it, a bullet whose text wraps to a second line
    reads as a separate bullet (the wrapped line falls back to the paragraph's left edge, not
    under the first line's text) -- this is what a hanging indent fixes.
    """
    pPr = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(width)))
    pPr.set("indent", str(int(-width)))


def text_units(s):
    """Width of one line in 'half-width' units: Hangul/Jamo count as 2, everything else as 1.

    Korean glyphs render roughly 2x the width of Latin ones in most UI fonts (Malgun Gothic
    included) -- summing raw character count under-predicts how many lines a Korean-heavy string
    will wrap to. Use this instead of len(s) when estimating wrap.
    """
    return sum(2 if ("가" <= ch <= "힣" or "ㄱ" <= ch <= "ㆎ") else 1 for ch in s)


def wrapped_row_count(lines, wrap_units=50):
    """Estimate total rendered rows for `lines` (each a paragraph/bullet string) once wrapped at
    `wrap_units` half-width units per line. Use to pre-size a card/textbox height before drawing it,
    so a two-line bullet doesn't get budgeted as one line and clip.

    wrap_units depends on the box width and font size in use -- calibrate per component (e.g. count
    units in a known-good rendered line at that width/size), don't assume 50 fits every layout.
    """
    return sum(1 + max(0, (text_units(ln) - 1) // wrap_units) for ln in lines)


def check_surface_leaks(prs, terms):
    """Scan every text frame in `prs` for any string in `terms` (editing traces, cross-document
    references, AI-tell phrasing -- e.g. "이전 버전", "다음과 같이", meta-commentary asides).

    Returns [(slide_no, term, snippet), ...] -- empty means clean.

    `terms` has NO default and is not itself shared: what counts as a "leak" is register- and
    project-specific (a phrase that's a tell in a lab-meeting deck may be normal prose in a teaching
    deck), and a broad shared list produces false positives (documented: "다음과 같다" flagged in a
    context where it was legitimate prose). Each project supplies its own list, informed by
    output_surface.md's register test, not this function.
    """
    hits = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for term in terms:
                if term in text:
                    hits.append((i + 1, term, text.strip()[:60]))
    return hits


def save_and_check(prs, path, leak_terms=None, sw=13.333, sh=7.5, tol=0.02):
    """Save + gate on surface leaks and boundary overflow in one call. Raises SystemExit (does not
    write a "saved" log line the caller can mistake for success) if either check fails -- a deck
    that fails either check should not reach the user un-flagged.

    Reuses `overflows()` per-slide (same tolerance as everywhere else in this module) rather than a
    second boundary check with different math, so there is one definition of "overflowing," not two
    that can disagree.
    """
    prs.save(path)
    leaks = check_surface_leaks(prs, leak_terms or [])
    overflow_hits = []
    for i, slide in enumerate(prs.slides):
        for name, r, b in overflows(slide, sw=sw, sh=sh, tol=tol):
            overflow_hits.append((i + 1, name, r, b))
    print("saved:", path, "slides=", len(prs.slides._sldIdLst))
    if leaks:
        print("[surface leaks]")
        for slide_no, term, snippet in leaks:
            print(f"  slide {slide_no}: '{term}' in \"{snippet}\"")
    if overflow_hits:
        print("[overflow]")
        for slide_no, name, r, b in overflow_hits:
            print(f"  slide {slide_no}: {name} right={r:.2f} bottom={b:.2f}")
    if leaks or overflow_hits:
        raise SystemExit(f"save_and_check failed -- leaks={len(leaks)}, overflow={len(overflow_hits)}. "
                          f"Fix and rebuild before delivering.")
    return prs
