# -*- coding: utf-8 -*-
"""poster_kit — palette-agnostic MECHANICS for large-format academic posters (python-pptx).

Same split as pptx_kit.py: this file has no colour/font opinions, callers pass those in. What lives
here is the fiddly, easy-to-get-wrong part of poster layout specifically — which is different enough
from a 16:9 slide deck (cm-scale canvas, CJK+Latin mixed text, a hard minimum legible font size read
from 1.5m away, figures whose *internal* labels must stay legible) that it doesn't belong in pptx_kit.

Extracted 2026-08-10 from `KWCS/teacher_violence_qol/R/03_poster.py`, whose 2-column English poster
was judged (by a co-author, informally, and by the user directly) clearly higher quality than a
contemporaneous 3-column poster built ad hoc for another project (`matching/brush_cog`). The
difference wasn't one trick — it was five habits every poster in the lab should now default to:

  1. **Two columns, not three.** Three columns forces small figures and short line lengths that read
     as cramped from viewing distance. Two columns give each figure real width and each line of text
     room to breathe. Use `two_col_grid()`.
  2. **One enforced minimum font size, no exceptions.** `kf()` asserts on every run — a poster is read
     from ~1.5m away, and "just this one caption a bit smaller" is how a poster ends up with a
     6pt-equivalent-effective-size line nobody proofread at that distance. Default 24pt; pass
     `min_pt=None` only for genuinely decorative micro-text (e.g. a copyright line), never for content.
  3. **Figures saved at layout size, not scaled down after.** `pic_cm()` locks width to the column and
     computes height from the real aspect ratio (PIL) — never pass a `max_h` that silently shrinks the
     image, because the point sizes *baked into the PNG* (axis labels, legends) shrink with it and fall
     below the 24pt floor invisibly. If a figure must be shorter, regenerate it at a shorter aspect
     ratio in the source plotting script — don't squeeze it here.
  4. **Narrative captions, not label captions.** `caption()` is happy to hold 2-3 sentences. "Figure 1.
     X vs Y" wastes the one place on a poster where a passerby who reads nothing else still gets the
     finding. Write what the figure *means*, the way an abstract's last sentence would.
  5. **One accent colour, used rarely.** Everything is BAR (section colour) or INK (body) or GRAY
     (captions/footnotes) except the single most important callout per section, which gets the accent
     colour — never four different tinted card-fill colours competing on one canvas. `bullets()` and
     `ptable()` both take an `accent` param for exactly this: mark ONE line or row, not several.

Pre-flight checklist before calling a poster "done" (beyond the usual pptx-editing skill steps):
  - No literal "TBD" / "to be confirmed" / placeholder text anywhere on the canvas — a poster is a
    finished artifact, not a draft; if something is genuinely unconfirmed, resolve it or drop it,
    don't ship the disclaimer (this happened for real: an author affiliation sat as "to be confirmed"
    on a rendered poster for two build cycles after the affiliation had already been confirmed
    elsewhere in the project).
  - `min_font_report()` returns empty (catches text created without going through `kf()`, e.g. hand-set
    table-cell runs).
  - `pptx_kit.overflows()` returns empty, using the poster's own sw_in/sh_in (NOT the 13.333x7.5 slide
    default — a poster canvas is tens of cm, always pass the real size).
  - Render via `render_pptx.py` and look at it at "arm's length" zoom, not full-screen — full-screen
    hides exactly the too-small-to-read-from-3-feet-away problem this kit exists to prevent.

Import path: lives in the pptx-editing skill's scripts/, alongside pptx_kit.py. Load it the same way
(see any *_poster.py in the lab for the `_load_pptx_kit`-style loader; add a second loader for this
file or extend that one to load both from the same directory).
"""
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

DEFAULT_MIN_PT = 24


def pt2cm(pt):
    return pt / 72 * 2.54


def kf(run, size, color, bold=False, italic=False, font="Arial", min_pt=DEFAULT_MIN_PT):
    """Set a run's font, enforcing the poster's minimum legible size.

    min_pt=None (or 0) opts a specific run out — reserve that for genuinely decorative text
    (a copyright line, a tiny logo caption), never for anything a reader is meant to learn from."""
    if min_pt:
        assert size >= min_pt, "font %.1fpt below poster minimum %dpt" % (size, min_pt)
    f = run.font
    f.size, f.color.rgb, f.bold, f.name, f.italic = Pt(size), color, bold, font, italic
    return run


def tb(sl, l, t, w, h, anchor=None):
    """A word-wrapped textbox at (l, t, w, h) in cm."""
    x = sl.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    x.text_frame.word_wrap = True
    if anchor:
        x.text_frame.vertical_anchor = anchor
    return x


def est_lines(txt, w_cm, size, char_w_factor=0.52, pad_cm=0.0):
    """Estimate wrapped line count for auto-sizing a textbox before drawing it.

    char_w_factor is the average glyph width as a fraction of the em (point size). This differs
    sharply by script — CALIBRATE, don't guess:
      - Latin/English (Arial):  ~0.52-0.60 (narrower glyphs, more chars per line)
      - Hangul:                 ~1.42       (roughly square glyphs, far fewer chars per line)
    Overestimating lines is the safe failure mode (extra whitespace, not a clipped line) — if
    unsure which way to round, round the factor UP.
    """
    usable = max(w_cm - pad_cm, 1)
    char_w = size / 72 * 2.54 * char_w_factor
    per_line = max(int(usable / char_w), 1)
    n = 0
    for seg in txt.split("\n"):
        n += max(1, -(-len(seg) // per_line))
    return n


def two_col_grid(page_w, margin, gutter, ncol=2):
    """Return (xs, colw) for an ncol-column grid spanning [margin, page_w - margin].

    Default the poster to 2 columns (see module docstring #1) — pass ncol=3 only when content
    genuinely cannot compress into two without a wall of tiny multi-panel figures; that was the
    actual failure mode this kit was extracted to fix, so treat 3 as the exception, not the default.
    """
    colw = (page_w - 2 * margin - (ncol - 1) * gutter) / ncol
    xs = [margin + i * (colw + gutter) for i in range(ncol)]
    return xs, colw


def sectitle(sl, txt, l, t, w, color, size=42, font="Arial", upper=True, min_pt=DEFAULT_MIN_PT):
    """Section header + underline rule. Returns the y (cm) where content should start next."""
    h = pt2cm(size) * 1.2
    x = tb(sl, l, t, w, h)
    kf(x.text_frame.paragraphs[0].add_run(), size, color, bold=True, font=font, min_pt=min_pt).text \
        = txt.upper() if upper else txt
    ry = t + h * 0.94
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(ry), Cm(w), Cm(0.2))
    r.fill.solid(); r.fill.fore_color.rgb = color; r.line.fill.background(); r.shadow.inherit = False
    return ry + 0.2 + 0.7


def _spans(txt):
    """'**bold** rest' -> [(text, is_bold), ...]. A line starting with a lone leading '**' with no
    matching close is treated as whole-line emphasis by the caller (bullets/ptable), not here."""
    out, i = [], 0
    for m in re.finditer(r"\*\*(.+?)\*\*", txt):
        if m.start() > i:
            out.append((txt[i:m.start()], False))
        out.append((m.group(1), True))
        i = m.end()
    if i < len(txt):
        out.append((txt[i:], False))
    return out or [(txt, False)]


def bullets(sl, items, l, t, w, ink, accent=None, size=27, gap=14, font="Arial", char_w_factor=0.52,
            mark="•  ", min_pt=DEFAULT_MIN_PT):
    """Auto-height bullet list. Inline `**bold**` spans; a whole item wrapped in a single leading
    `**...` (i.e. starts with '**' and has no closing pair elsewhere) is rendered entirely in
    `accent` colour+bold — use this for AT MOST one bullet per section (module docstring #5), not as
    a general emphasis tool. Returns the total height (cm) consumed."""
    plain = [s.replace("**", "") for s in items]
    total = sum(est_lines(s, w - 1.2, size, char_w_factor) for s in plain) * pt2cm(size) * 1.2 \
        + len(items) * pt2cm(gap)
    x = tb(sl, l, t, w, total + 1.0)
    tf = x.text_frame
    for i, s in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i:
            p.space_before = Pt(gap)
        whole = s.startswith("**") and s.count("**") == 1
        s2 = s[2:] if whole else s
        whole_color = accent if (whole and accent) else ink
        kf(p.add_run(), size, whole_color, bold=whole, font=font, min_pt=min_pt).text = mark
        for seg, bd in _spans(s2):
            col = accent if ((whole or bd) and accent) else ink
            kf(p.add_run(), size, col, bold=(whole or bd), font=font, min_pt=min_pt).text = seg
    return total + 1.0


def caption(sl, txt, l, t, w, color, accent=None, size=24, font="Arial", char_w_factor=0.52,
            min_pt=DEFAULT_MIN_PT):
    """A narrative caption (module docstring #4 — write what the figure MEANS). Returns height (cm).

    Supports the same inline `**bold**` markup as bullets() (rendered bold, in `accent` colour if
    given, else just bold in `color`) — pass `accent` if the caption emphasises one clause. A caption
    with no `**` in it renders as a single plain run, unchanged from before this parameter existed.
    """
    plain = txt.replace("**", "")
    h = est_lines(plain, w, size, char_w_factor) * pt2cm(size) * 1.2 + 0.5
    x = tb(sl, l, t, w, h)
    p = x.text_frame.paragraphs[0]
    for seg, bd in _spans(txt):
        col = accent if (bd and accent) else color
        kf(p.add_run(), size, col, bold=bd, font=font, min_pt=min_pt).text = seg
    return h


def ptable(sl, rows, l, t, w, widths, header_color, alt_color, white, ink, accent=None, size=24,
           rowh=1.9, title=None, font="Arial", min_pt=DEFAULT_MIN_PT, char_w_factor=0.55):
    """Poster table. `**cell` bolds that cell; a title (rendered as a merged first row, matching the
    convention seen in published posters) goes in `title`, not as a separate textbox above the table
    — keeps the table a single self-contained object when the layout height is recalculated.

    `rowh` is a FLOOR, not the assumed height of every row. Each row's real height is computed from
    the wrapped line count of its longest cell (est_lines(), the same estimator bullets()/caption()
    already use for prose) — a row only gets taller than rowh if some cell in it actually needs to
    wrap at that column's width. The returned height is the true sum, not `rowh * nrow`.

    Why this matters (2026-08-14, brush_cog poster, hit 3 times independently in one build before
    this fix existed): a plain `rowh * nrow` silently under-predicts any row whose cell wraps — a
    long row label ("After-lunch brushing" in a narrow first column), a header cell ("CHS
    (N=79,826)" in a column sized for short data values), or a long table `title`. PowerPoint then
    auto-grows JUST that row at render time (this is real, static geometry checks can't see it —
    see deck_render_audit.py's docstring), so the table's true rendered height exceeds what this
    function told the caller, and whatever the caller places next (almost always a caption)
    overlaps the table's actual last row. No amount of nudging the caption's y-offset fixes this
    reliably, because the mismatch scales with content, not a fixed pixel amount — the fix has to
    be that the returned height is right in the first place.
    """
    ncol = len(widths)
    col_w = [ww * (w / sum(widths)) for ww in widths]
    line_h = pt2cm(size) * 1.2

    def cell_lines(val, cw):
        plain = val[2:] if val.startswith("**") else val
        return est_lines(str(plain), cw, size, char_w_factor, pad_cm=0.4)

    def row_height(rr):
        return max(rowh, max(cell_lines(v, col_w[j]) for j, v in enumerate(rr)) * line_h + 0.3)

    row_heights = []
    if title:
        title_lines = est_lines(title, w, size + 2, char_w_factor, pad_cm=0.35)
        row_heights.append(max(rowh, title_lines * pt2cm(size + 2) * 1.2 + 0.3))
    row_heights += [row_height(rr) for rr in rows]
    total_h = sum(row_heights)

    nrow = len(rows) + (1 if title else 0)
    gt = sl.shapes.add_table(nrow, ncol, Cm(l), Cm(t), Cm(w), Cm(total_h))
    tblx = gt.table
    for j, ww in enumerate(widths):
        tblx.columns[j].width = Cm(col_w[j])
    off = 0
    if title:
        tblx.rows[0].height = Cm(row_heights[0])
        c0 = tblx.cell(0, 0)
        c0.merge(tblx.cell(0, ncol - 1))
        c0.fill.solid(); c0.fill.fore_color.rgb = white
        c0.margin_left = c0.margin_right = Cm(0.15)
        kf(c0.text_frame.paragraphs[0].add_run(), size + 2, header_color, bold=True, font=font,
           min_pt=min_pt).text = title
        off = 1
    for i, rr in enumerate(rows):
        tblx.rows[i + off].height = Cm(row_heights[i + off])
        head = (i == 0)
        for j, val in enumerate(rr):
            c = tblx.cell(i + off, j)
            c.margin_left = c.margin_right = Cm(0.18)
            c.margin_top = c.margin_bottom = Cm(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            c.fill.fore_color.rgb = header_color if head else (alt_color if (i % 2 == 0) else white)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            bold = head or val.startswith("**")
            val2 = val[2:] if val.startswith("**") else val
            col = white if head else (accent if (val.startswith("**") and accent) else ink)
            kf(p.add_run(), size, col, bold=bold, font=font, min_pt=min_pt).text = val2
    return total_h


def stack_box_height(w, lines, size=24, char_w_factor=0.55, pad_cm=0.5):
    """The height stack_box() would use for this (w, lines) — call this FIRST when two or more
    boxes must share one height (e.g. two lanes of a flowchart row): compute each side's natural
    height, take the max, then pass it as `min_h` to every stack_box() call in that row so they
    align without any box being under-sized for its own content."""
    line_h = pt2cm(size) * 1.2
    n_lines = sum(est_lines(txt, w - pad_cm, size, char_w_factor) for txt, _, _ in lines)
    return n_lines * line_h + pad_cm


def stack_box(sl, l, t, w, lines, fill, border=None, size=24, font="Arial",
              min_pt=DEFAULT_MIN_PT, char_w_factor=0.55, pad_cm=0.5, align=PP_ALIGN.CENTER,
              min_h=0.0):
    """A rounded-rect box auto-sized to fit `lines` (a list of (text, bold, color) tuples, one per
    line, centred and stacked top-to-bottom) — height computed from est_lines() at the box's own
    width, the same pattern ptable() now uses for cells. Returns (shape, height_cm). For several
    boxes that must share one height, see stack_box_height() first.

    Extracted 2026-08-14 from a project-local flowchart box() that was hand-sized with a guessed
    constant three separate times in the same build before someone finally computed it from the
    real line count — this is that computation, done once, shared, so the next STROBE/flow-diagram
    box in any project starts from a box that actually fits its own text instead of reinventing
    (and re-breaking) the same helper. Multi-line wrapping WITHIN one logical line is estimated via
    est_lines(); this does not merge/re-wrap the list itself, so pass already-short phrases (one
    idea per list item), not paragraphs.
    """
    h = max(min_h, stack_box_height(w, lines, size, char_w_factor, pad_cm))
    r = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    r.adjustments[0] = 0.03
    if border:
        r.line.color.rgb = border; r.line.width = Pt(1.5)
    else:
        r.line.fill.background()
    r.shadow.inherit = False
    b = tb(sl, l + pad_cm / 2, t, w - pad_cm, h, anchor=MSO_ANCHOR.MIDDLE)
    for i, (txt, bold, color) in enumerate(lines):
        p = b.text_frame.paragraphs[0] if i == 0 else b.text_frame.add_paragraph()
        p.alignment = align
        kf(p.add_run(), size, color, bold=bold, font=font, min_pt=min_pt).text = txt
    return r, h


def pic_cm(sl, path, l, t, w_cm):
    """Insert an image at exactly `w_cm` wide, height from its real aspect ratio (PIL).

    Deliberately has NO max_h parameter (module docstring #3) — a poster figure's width should be set
    by the column it lives in, and its height should follow from that. If a figure is too tall for the
    space left on the page, the fix is regenerating the figure at a different aspect ratio (in its own
    plotting script) or reordering sections, not shrinking it here — shrinking silently drops any
    in-figure text below the poster's minimum legible size.
    """
    from PIL import Image
    iw, ih = Image.open(path).size
    h_cm = w_cm * ih / iw
    sl.shapes.add_picture(path, Cm(l), Cm(t), Cm(w_cm), Cm(h_cm))
    return h_cm


def min_font_report(slide, min_pt=DEFAULT_MIN_PT):
    """Mechanical pre-flight check: walk every run on the slide (textboxes AND table cells) and
    return [(shape_name, text_snippet, size_pt), ...] for anything under min_pt.

    Catches text created without going through kf() — e.g. a table cell font set by hand, or a run
    whose size was computed and happens to round under the floor. kf()'s assert only protects text
    that actually goes through it; this is the safety net for text that didn't. Empty return = clean.

    Narrow by design: this only re-checks the [RUN] case (a literal font size on a run). It does NOT
    catch an image placed smaller than native — baked-in labels shrink with it, 11pt drawn at 0.7x
    placement renders at 7.7pt — or an undersized cex=/size=/fontsize= literal in the R/Python script
    that generated a figure. This function has no dependency outside this file, so a standalone
    install of the public skill only gets this narrower check. If the lab's shared tooling is also
    available, run `agent/tools/audit_font_sizes.py` on the built .pptx for the fuller three-pronged
    check (RUN + SCALE + SRC) — it is not imported here, so the two can drift; if the definition of
    "too small" changes in one, check whether it should change in the other."""
    hits = []
    for shp in slide.shapes:
        runs = []
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                runs.extend(p.runs)
        if shp.has_table:
            for row in shp.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        runs.extend(p.runs)
        for r in runs:
            sz = r.font.size
            if sz is not None and sz.pt < min_pt and r.text.strip():
                hits.append((shp.name, r.text.strip()[:40], round(sz.pt, 1)))
    return hits
