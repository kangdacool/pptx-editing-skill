# -*- coding: utf-8 -*-
"""audit_surface_text.py FILE.pptx [--extra-pattern RE ...]

배포 직전 «표면 텍스트» 감사. audit_text_fit.py가 «글자가 넘치는가»를 본다면,
이 스크립트는 «있으면 안 되는 말이 남았는가»를 본다. 둘은 다른 실패를 잡으므로 둘 다 돌린다.

왜 필요한가 (2026-08-10, teacher_violence_qol):
  같은 세션에서 사용자가 "쓸데없는 서술이 많다"고 두 번 지적했고, 지운 뒤에도 즉석 grep을
  돌리자 4건이 더 나왔다. 사람 눈으로 훑는 방식은 반복해서 실패한다 -- 기계로 훑어야 한다.

무엇을 잡나 (기본 5종 -- 앞 4종은 «편집 흔적», 마지막은 «수사»):
  meta        편집 해명·검정 의미 해설·자기지시적 회고 ("justified", "not just",
              "are not interpreted", "already unreliable" -- 2026-08-16 brush_cog
              Table 5 캡션에서 실제로 잡힌 사례: 독자가 처음 보는 표인데 "already"로
              마치 앞서 논의한 것처럼 서술)
  nav         내비게이션 안내 ("see Table 3", "column 2", "decomposed in")
  restate     재진술 신호 ("I.e.,", "in other words")
  provenance  내부 파일명·파이프라인 잔재 (".csv", "Source: TableX_...")
  rhetoric    ⭐ 슬라이드의 «AI가 만든 티»는 내용이 아니라 수사다 -- 카드 제목이 그 상자의
              «수사적 기능»("The hook"·"Key point"), 추정치 뒤 화살표, 그리고 덱 전체 문맥으로
              판정하는 대문자 강조(caps_emphasis: 같은 낱말이 ALLCAPS와 소문자로 둘 다 나오면
              약어가 아니라 강조다 -- 그 대조 하나로 NHANES와 DEPLETES가 갈린다).
              판정은 금칙어가 아니라 «화자»다: 동료에게 연구를 말하는 저자인가, 논문을 파는
              편집자인가. 경구형 대구·`≠`/`≈`는 일부러 뺐다(각각 판단·register 의존).

번호 drift(파이프라인 번호가 산출물에 남은 것)는 --max-fig/--max-tab로 잡는다:
  산출물에 Figure가 3개뿐인데 "Figure 9"가 캡션에 남아 있으면 청중에게 없는 그림을 가리킨다.

exit 1 on any hit -- 빌드 파이프라인에서 게이트로 쓸 수 있다.

이 파일의 DEFAULT는 **다른 감사 도구들과 다른 철학**이다 -- `pptx_kit.check_surface_leaks`와
`agent/tools/surface_leak_scan.py`는 의도적으로 기본 목록이 없다("register마다 안전한 문구가
다르므로 공유 기본값은 어느 한쪽에서 오탐·누락을 낸다"는 것이 그 두 도구의 명시된 설계 원칙).
이 파일이 DEFAULT를 가진 이유는 대상이 다르기 때문이다 -- 아래 4종은 영문 학술 산출물 표면에서
**register와 거의 무관하게 늘 문제인** 기계적 패턴("justified", "see Table 3", "I.e.,", 내부
파일명)만 골랐다. 문서별로 안전성이 갈리는 표현(예: "다음과 같다"가 정당한 산문인 문서도 있다)은
DEFAULT에 넣지 않았고, 그런 표현은 `--extra-pattern`으로 이 문서에 한해 추가한다 -- 그러면 이
스크립트도 사실상 "기계적 베이스라인 + caller-supplied 추가"가 되어 다른 두 도구와 같은 원칙을
따르게 된다. 셋의 관계: **이 파일**은 pptx 산출물 전용 베이스라인, **surface_leak_scan.py**는
어떤 텍스트 파일에도 쓰는 register-specific 전용(기본값 없음), **check_surface_leaks**는 그
철학을 파이썬 객체 저장 시점에 인라인으로 거는 게이트다.
"""
import re
import sys
import argparse

DEFAULT = {
    "meta": r"\bjustified\b|\bnot just\b|are not interpreted|separation artifact|"
            r"\brobustness\b|sensitivity analys|\bverified\b|directly testing|"
            r"\bnotably\b|\bimportantly\b|\bhonestly\b|it is worth noting|"
            r"\balready\s+(unreliable|noted|established|known|discussed|mentioned|shown|seen)\b",
    "nav": r"see (Table|Figure|column|panel)|decomposed in|\bcolumn \d|as shown (above|below)|"
           r"refer to (Table|Figure)",
    "restate": r"\bI\.e\.|\bi\.e\.,|in other words|that is to say|this means that",
    "provenance": r"\.csv\b|\.rds\b|\.xlsx\b|Source: [A-Z][A-Za-z0-9_]*_|output/tables",
    # 2026-08-20 신설. 위 넷은 «편집 흔적»(내가 편집자로서 남긴 말)이고, 이건 «수사»다 --
    # 문장이 전부 참인데도 «생성물»로 읽히게 만드는 형태. 규칙은 [[ppt-rules]]에 2026-08-03부터
    # 있었고 grep 한 줄까지 적혀 있었는데 도구가 없어서 아무도 안 돌렸다.
    #   ① 카드 제목이 그 상자의 «수사적 기능»  -- `The hook`은 논문을 파는 내부 용어다.
    #      상자에 «든 것»을 제목으로 써야 한다. 닫힌 집합이라 register 무관하게 안전하다.
    #   ② 기호로 쓴 말 -- 추정치 뒤 화살표. 글로 쓴다.
    # ⚠️ 여기 «넣지 않은» 것과 그 이유:
    #   · `≠`/`≈` -- 통계 덱에서는 정당한 수학 기호다. register 의존이라 DEFAULT에 안 맞는다.
    #   · 경구형 대구("Rigorous analysis is the strength; the design is the ceiling") -- 판단이다.
    #     정규식으로 박으면 정상적인 대조 문장을 죽인다.
    #   · 문장 속 대문자 강조 -- 약어(NHANES·STROBE·TMLE)와 구분이 안 된다. 정규식 대신
    #     아래 caps_emphasis()가 «덱 전체» 문맥으로 판정한다.
    "rhetoric": r"^\s*(The\s+(hook|consequence|takeaway|upshot|punchline)|"
                r"Key\s+point|Organizing\s+principle|Why\s+this\s+matters)\s*[:.]?\s*$|"
                r"[↑↓]",
}


def caps_emphasis(paras):
    """대문자 강조를 «덱 전체» 문맥으로 판정한다 -- 정규식만으로는 약어와 구분이 안 되기 때문.

    신호: 같은 낱말이 이 덱 안에서 ALLCAPS로도 소문자로도 나온다면 그건 약어가 아니라 강조다.
    NHANES·STROBE·TMLE 같은 진짜 약어는 소문자로 등장하는 일이 없다. 이 한 가지 대조가
    «DEPLETES»(강조)와 «NHANES»(약어)를 갈라준다 -- 화이트리스트를 손으로 관리할 필요가 없다.

    ⚠️ 그 대조만으로는 부족하다: 표 헤더는 설계상 ALLCAPS인 경우가 많아 «MEAN» 헤더 + 본문
    «mean»이 오탐이 된다(셀프테스트를 쓰다 발견). 그래서 «문장 안»일 것을 함께 요구한다 --
    강조는 문장 중간에 박히고, 헤더는 그 칸 전체다. 조건: 그 문단에 낱말이 4개 이상이고
    소문자 낱말이 함께 있을 것.
    """
    lower, caps = set(), {}
    for si, t in paras:
        words = re.findall(r"\b[A-Za-z]{4,}\b", t)
        in_sentence = len(t.split()) >= 4 and any(w.islower() for w in words)
        for w in words:
            if w.isupper():
                if in_sentence:
                    caps.setdefault(w.lower(), (si, w))
            elif w.islower():
                lower.add(w.lower())
    return [("rhetoric-caps", si, w, "이 덱에 소문자 %r도 있다 -- 약어가 아니라 강조다" % w.lower())
            for k, (si, w) in sorted(caps.items()) if k in lower]


def paragraphs(path):
    from pptx import Presentation
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        yield si, t
            if getattr(sh, "has_table", False):
                for r in sh.table.rows:
                    for c in r.cells:
                        t = c.text.strip()
                        if t:
                            yield si, t


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--extra-pattern", action="append", default=[],
                    help="추가 정규식 (프로젝트별 금지어). 여러 번 지정 가능.")
    ap.add_argument("--max-fig", type=int, default=None,
                    help="산출물의 Figure 최대 번호. 초과 참조를 번호 drift로 신고.")
    ap.add_argument("--max-tab", type=int, default=None,
                    help="산출물의 Table 최대 번호.")
    ap.add_argument("--skip", action="append", default=[], choices=sorted(DEFAULT),
                    help="이 register에선 정상인 기본 범주를 끈다. 예: 내부 회의 덱은 "
                         "'출처: xxx.csv' 표기가 **바람직하므로** --skip provenance.")
    a = ap.parse_args()

    pats = {k: re.compile(v, re.I) for k, v in DEFAULT.items() if k not in a.skip}
    if a.skip:
        print("(register상 제외한 범주: %s)" % ", ".join(sorted(set(a.skip))))
    for i, ex in enumerate(a.extra_pattern):
        pats["custom%d" % (i + 1)] = re.compile(ex, re.I)

    paras = list(paragraphs(a.file))
    hits = []
    for si, t in paras:
        for name, pat in pats.items():
            m = pat.search(t)
            if m:
                hits.append((name, si, m.group(0), t))

    def drift(kind, cap):
        out = []
        for si, t in paras:
            for m in re.finditer(r"\b%s\s*(\d+)" % kind, t, re.I):
                if int(m.group(1)) > cap:
                    out.append(("number-drift", si, m.group(0), t))
        return out

    if a.max_fig is not None:
        hits += drift("Figure", a.max_fig)
    if a.max_tab is not None:
        hits += drift("Table", a.max_tab)
    if "rhetoric" not in a.skip:
        hits += caps_emphasis(paras)

    print("%s -- %d text blocks scanned" % (a.file, len(paras)))
    if not hits:
        print("  OK  surface-text audit clean")
        return 0
    print("  %d hit(s):" % len(hits))
    for name, si, frag, t in hits:
        print("\n  [%s] slide %d  matched %r" % (name, si, frag))
        print("      " + (t[:220] + ("..." if len(t) > 220 else "")))
    return 1


if __name__ == "__main__":
    sys.exit(main())
