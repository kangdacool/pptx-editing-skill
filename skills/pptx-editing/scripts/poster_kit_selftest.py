# -*- coding: utf-8 -*-
"""poster_kit_selftest.py — prove poster_kit's mechanics without a real poster project.
    python poster_kit_selftest.py   # prints PASS / raises on failure
"""
import sys, io, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
if hasattr(sys.stdout, "buffer"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
import poster_kit as pk

INK = RGBColor(0x1A, 0x1A, 0x1A)
BAR = RGBColor(0x1F, 0x4E, 0x79)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xEC, 0xF2, 0xF8)
RED = RGBColor(0xB0, 0x30, 0x2D)


def blank_slide(pw_cm=90.0, ph_cm=120.0):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(pw_cm), Cm(ph_cm)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_min_pt_enforced():
    _, sl = blank_slide()
    x = pk.tb(sl, 1, 1, 10, 2)
    raised = False
    try:
        pk.kf(x.text_frame.paragraphs[0].add_run(), 18, INK)  # below DEFAULT_MIN_PT=24
    except AssertionError:
        raised = True
    assert raised, "kf() must reject sub-minimum font sizes by default"
    pk.kf(x.text_frame.paragraphs[0].add_run(), 24, INK)  # exactly at floor: OK
    pk.kf(x.text_frame.paragraphs[0].add_run(), 10, INK, min_pt=None)  # explicit opt-out: OK


def test_two_col_grid():
    xs, colw = pk.two_col_grid(page_w=90.0, margin=3.0, gutter=2.5)
    assert len(xs) == 2
    assert abs(xs[0] - 3.0) < 1e-9
    assert abs(colw - 40.75) < 1e-9
    assert abs(xs[1] - (xs[0] + colw + 2.5)) < 1e-9


def test_est_lines_direction():
    # Hangul (large char_w_factor) should estimate MORE lines than Latin for the same string/width.
    txt = "x" * 40
    latin = pk.est_lines(txt, w_cm=20, size=24, char_w_factor=0.55)
    hangul = pk.est_lines(txt, w_cm=20, size=24, char_w_factor=1.42)
    assert hangul >= latin


def test_bullets_sectitle_caption_ptable_render():
    _, sl = blank_slide()
    y = pk.sectitle(sl, "Results", 3, 3, 40, BAR)
    y += pk.bullets(sl, ["plain bullet", "**whole-line accent bullet", "mixed **bold** span"],
                     3, y, 40, INK, accent=RED)
    y += pk.caption(sl, "This figure shows the finding, narratively, not just a label.",
                     3, y, 40, GRAY)
    y += pk.ptable(sl, [["A", "B"], ["1", "**2"]], 3, y, 40, [20, 20],
                    header_color=BAR, alt_color=PALE, white=WHITE, ink=INK, accent=RED,
                    title="Table 1. Demo")
    assert y > 3  # something was drawn


def test_ptable_row_grows_for_wrapped_cell():
    """2026-08-14 regression: a row whose cell text wraps at the column's real width must get a
    taller row (and a taller total returned height) than a row of short cells at the same rowh —
    otherwise the caller places the next element (usually a caption) using an undershot height and
    it lands on top of the table's real last row. Hit 3x independently in the same poster build
    (a row label, a table header cell, a wrapped title) before ptable() computed this itself."""
    _, sl = blank_slide()
    narrow_col = [6, 6]  # cm - narrow enough that a long phrase must wrap
    short_rows = [["Hdr", "Hdr"], ["A", "B"]]
    h_short = pk.ptable(sl, short_rows, 3, 3, sum(narrow_col), narrow_col,
                         header_color=BAR, alt_color=PALE, white=WHITE, ink=INK, rowh=1.5)
    long_rows = [["Hdr", "Hdr"], ["This label is long enough that it must wrap here", "B"]]
    h_long = pk.ptable(sl, long_rows, 3, 30, sum(narrow_col), narrow_col,
                        header_color=BAR, alt_color=PALE, white=WHITE, ink=INK, rowh=1.5)
    assert h_long > h_short, (
        "a wrapped-cell table must report a taller total height than an all-short-cell table "
        "at the same rowh — otherwise ptable() is still under-predicting wrap")
    # the actual row height set on the shape must also reflect the wrap, not just the return value
    tbl_long = sl.shapes[-1].table
    assert tbl_long.rows[1].height > Cm(1.5 + 0.01), (
        "the wrapped data row's own .height must exceed the rowh floor")


def test_ptable_wrapped_title_matches_returned_height():
    """A title too long for one line must grow the title row AND the returned total height by the
    same amount PowerPoint will actually render, not just the ad-hoc '+1 row' guess a caller-side
    correction used before this was computed inside ptable() itself."""
    _, sl = blank_slide()
    widths = [8, 8]
    long_title = "This title is deliberately long enough that it will not fit on one line at this width"
    h = pk.ptable(sl, [["A", "B"], ["1", "2"]], 3, 3, sum(widths), widths,
                  header_color=BAR, alt_color=PALE, white=WHITE, ink=INK, title=long_title)
    tbl = sl.shapes[-1].table
    # returned height must equal the sum of the ACTUAL row heights ptable() set - if these
    # disagree, whatever the caller places next will be mispositioned relative to the real table.
    actual_total_cm = sum(row.height for row in tbl.rows) / 360000.0  # EMU -> cm
    assert abs(h - actual_total_cm) < 0.05, (
        "ptable()'s returned height must match the sum of the row heights it actually set")


def test_stack_box_sized_for_its_own_content():
    """stack_box() must grow for more lines - a box for 3 lines must be taller than a box for 1
    line at the same width/font, and must not raise the kf() min-size assertion."""
    _, sl = blank_slide()
    _, h1 = pk.stack_box(sl, 3, 3, 15, [("One line", True, INK)], fill=WHITE)
    _, h3 = pk.stack_box(sl, 3, 20, 15,
                          [("Line one", True, INK), ("Line two", False, GRAY),
                           ("Line three", False, GRAY)], fill=WHITE)
    assert h3 > h1, "a 3-line stack_box must be taller than a 1-line one"


def test_min_font_report_catches_bypass():
    _, sl = blank_slide()
    # a run created WITHOUT poster_kit.kf(), i.e. bypassing the floor
    x = pk.tb(sl, 1, 1, 10, 2)
    r = x.text_frame.paragraphs[0].add_run()
    r.text = "too small"
    from pptx.util import Pt
    r.font.size = Pt(12)
    hits = pk.min_font_report(sl, min_pt=24)
    assert any("too small" in h[1] for h in hits), "min_font_report must catch non-kf() runs"


if __name__ == "__main__":
    for name, fn in list(globals().items()):
        if name.startswith("test_"):
            fn()
            print("PASS:", name)
    print("ALL PASS")
