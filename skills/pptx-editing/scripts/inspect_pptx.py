# -*- coding: utf-8 -*-
"""inspect_pptx.py — per-slide structure + speaker-notes + boundary-overflow dump.

Run this FIRST on any deck (especially one you did not build) before editing it.
    python inspect_pptx.py FILE.pptx [--text]

Prints, per slide: index, layout, shape/table/picture counts, speaker-note length + preview, and any
shape whose right/bottom leaves the slide (the thing python-pptx will not warn you about). --text also
dumps each text box's content (truncated)."""
import sys, io, argparse
if hasattr(sys.stdout, "buffer"):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0

def note_text(slide):
    if not slide.has_notes_slide:
        return None
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        return None
    t = tf.text or ""
    return t if t.strip() else None

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("file")
    ap.add_argument("--text", action="store_true", help="dump text-box contents (truncated)")
    a = ap.parse_args()
    prs = Presentation(a.file)
    sw, sh = prs.slide_width / EMU, prs.slide_height / EMU
    print(f"{a.file} :: {len(prs.slides)} slides :: {sw:.2f} x {sh:.2f} in")
    n_notes = 0
    for i, s in enumerate(prs.slides, 1):
        shapes = list(s.shapes)
        tables = sum(1 for x in shapes if x.has_table)
        pics = sum(1 for x in shapes if x.shape_type == 13)  # PICTURE
        over = []
        for x in shapes:
            try:
                r = (x.left + x.width) / EMU; b = (x.top + x.height) / EMU
            except TypeError:
                continue
            if r > sw + 0.02 or b > sh + 0.02:
                over.append(f"{x.name}(r={r:.2f},b={b:.2f})")
        layout = s.slide_layout.name
        nt = note_text(s)
        if nt: n_notes += 1
        print(f"\n[{i:>2}] layout={layout!r}  shapes={len(shapes)} tables={tables} pics={pics}"
              f"{'  NOTE(%d)' % len(nt) if nt else ''}"
              f"{'  OVERFLOW: ' + ', '.join(over) if over else ''}")
        if nt:
            print("      note: " + nt[:120].replace("\n", " / "))
        if a.text:
            for x in shapes:
                if x.has_text_frame and x.text_frame.text.strip():
                    txt = x.text_frame.text.strip().replace("\n", " / ")
                    print(f"      txt: {txt[:110]}")
    print(f"\n== {n_notes} slide(s) carry speaker notes ==")

if __name__ == "__main__":
    main()
